#!/usr/bin/env python3
"""Scaffold a guide-first deck workspace."""

from __future__ import annotations

import argparse
import sys
import textwrap
from pathlib import Path


def slide_id(index: int, total: int) -> str:
    if index == 1:
        suffix = "cover"
    elif index == total:
        suffix = "closing"
    elif index == total - 1 and total >= 3:
        suffix = "summary"
    else:
        suffix = f"content-{index:02d}"
    return f"s{index:02d}-{suffix}"


def english_guide(title: str, audience: str, duration: int, slides: int, author: str, generator_name: str, backend_name: str) -> str:
    sections = []
    for idx in range(1, slides + 1):
        sid = slide_id(idx, slides)
        sections.append(
            textwrap.dedent(
                f"""\
                ### {sid}

                - Target Duration: `TODO`
                - Cumulative Time: `TODO`
                - Title: `TODO`
                - Subtitle: `TODO`
                - Message: `TODO`
                - Visible Text:
                  - `TODO`
                - Visuals:
                  - `TODO`
                - Formula Requirements:
                  - `None` or `TODO`
                - Symbol Explanations:
                  - `None` or `TODO`
                - Speaker Notes:
                  - `TODO`
                - Notes Supplement:
                  - `None` or `TODO`
                - Implementation Notes:
                  - `TODO`
                - Generator-Ready Instructions:
                  - `TODO`
                - Acceptance Checks:
                  - `TODO`
                """
            ).rstrip()
        )

    header = textwrap.dedent(
        f"""\
        # PPT_GUIDE

        ## Deck Brief

        - Title: `{title}`
        - Goal: `TODO`
        - Audience: `{audience}`
        - Duration: `{duration} minutes`
        - Slide Budget: `{slides}`
        - Presenter: `{author}`
        - Aspect Ratio: `16:9`
        - Tooling Backend: `{backend_name}`
        - Notes Policy: `script-only` or `script-plus-supplement`
        - Audience Language: `TODO`
        - Source Materials:
          - `TODO`
        - Style Constraints:
          - `TODO`
        - Acceptance Criteria:
          - `Rendered deck has no obvious overflow, collision, or clipping`
          - `The final deck matches the slide messages below`
          - `Internal slide ids such as s01-cover stay in the guide, code, and review notes instead of visible slide text unless explicitly requested`

        ## Document Purpose And Usage

        - This guide is the source of truth for deck intent.
        - Keep this workspace inside a durable user repo, not in a transient temp directory.
        - Update this guide first for narrative, scope, timing, and slide-order changes.
        - Update `{generator_name}` for layout and implementation changes.
        - Re-render after every meaningful edit.
        - If the final deck must embed speaker notes, treat this guide as the notes authority.

        ## Timing Plan

        - `s01-cover`: `TODO`
        - `s02-content-02`: `TODO`

        ## Global Production Principles

        - Keep visible slide text audience-facing.
        - Keep presenter-facing explanation in speaker notes.
        - Keep important formulas on-screen when they carry the slide's core argument, and explain the key symbols.
        - Keep one main message per slide.
        - Keep stable slide ids in the guide, generator, comments, and review notes, not on visible slides unless explicitly requested.
        - Keep summary, takeaway, and closing-adjacent visible text in the audience's working language unless the user explicitly wants otherwise.
        - Make each slide specific enough that the generator can implement it without guessing.

        ## Slide Plan
        """
    ).rstrip()
    return header + "\n\n" + "\n\n".join(sections) + "\n"


def chinese_guide(title: str, audience: str, duration: int, slides: int, author: str, generator_name: str, backend_name: str) -> str:
    sections = []
    for idx in range(1, slides + 1):
        sid = slide_id(idx, slides)
        sections.append(
            textwrap.dedent(
                f"""\
                ### {sid}

                - 目标时长：`TODO`
                - 累计时间：`TODO`
                - 页面标题：`TODO`
                - 页面副标题（可选）：`TODO`
                - 页面核心信息：`TODO`
                - 页面上直接放的字：
                  - `TODO`
                - 图片 / 图表要求：
                  - `TODO`
                - 公式要求（可选）：
                  - `无` 或 `TODO`
                - 符号解释（可选）：
                  - `无` 或 `TODO`
                - 可直接念的完整台词：
                  - `TODO`
                - notes 补充（可选）：
                  - `无` 或 `TODO`
                - 实现备注：
                  - `TODO`
                - 可供 generator 直接照办的说明：
                  - `TODO`
                - 检查重点：
                  - `TODO`
                """
            ).rstrip()
        )

    header = textwrap.dedent(
        f"""\
        # PPT_GUIDE

        ## 默认假设

        - 标题：`{title}`
        - 听众：`{audience}`
        - 总时长：`{duration} 分钟`
        - 页数：`{slides}`
        - 汇报人：`{author}`
        - 画布：`16:9`
        - 生成后端：`{backend_name}`
        - notes 口径：`仅等于可直接念台词` 或 `可直接念台词 + notes 补充`
        - 观众语言：`TODO`

        ## 文档定位与使用方式

        - 这份 guide 是 deck 内容与结构的上游规范。
        - 这套工作区应放在用户的持久化 repo 里，而不是临时目录里。
        - 改主线、结构、页数、讲稿，先改这份 guide。
        - 改布局、裁图、字号、配色、遮挡，先改 `{generator_name}`。
        - 每次实质修改后都重新渲染检查。
        - 如果最终 deck 需要内嵌 notes，这份 guide 就是 notes 的上游真值。
        - `s01-cover` 这类 slide id 仅用于 guide、代码、review 追踪，不应直接放到观众可见页面上，除非用户明确要求。

        ## 总时长分配

        - `s01-cover`：`TODO`
        - `s02-content-02`：`TODO`

        ## 制作总原则

        - 所有可见文字都应面向观众，而不是面向讲解者。
        - 讲稿与 note 可以保留讲解者视角，但不要把讲解提示写进可见区。
        - 重要公式如果承载这页的核心论点，应真正上屏，并明确哪些符号要解释给观众。
        - 每页只承担一个主结论或主问题。
        - slide id 是制作期元数据，不是观众内容。
        - 总结页、takeaway 页和封底前一页默认应保持观众语言一致，长段异语言解释优先放 notes。
        - 页面说明必须具体到 generator 可以直接照办。

        ## 逐页规划
        """
    ).rstrip()
    return header + "\n\n" + "\n\n".join(sections) + "\n"


def js_template(title: str, author: str) -> str:
    return textwrap.dedent(
        f"""\
        const pptxgen = require("pptxgenjs");

        function buildCover(pptx) {{
          // Keep stable slide ids in function names, comments, and review notes.
          // Do not place ids like s01-cover on the visible slide unless explicitly requested.
          // If a slide needs an important formula, keep it visible and explain symbols instead of hiding it only in notes.
          const slide = pptx.addSlide();
          slide.background = {{ color: "F7F5F1" }};
          slide.addText("{title}", {{
            x: 0.7,
            y: 0.8,
            w: 10.5,
            h: 0.6,
            fontFace: "Aptos Display",
            fontSize: 24,
            bold: true,
            color: "1F2937",
            margin: 0,
          }});
          slide.addText("Replace this placeholder with audience-facing content from PPT_GUIDE.md.", {{
            x: 0.75,
            y: 1.6,
            w: 11.4,
            h: 0.8,
            fontFace: "Aptos",
            fontSize: 14,
            color: "4B5563",
            margin: 0,
          }});
          slide.addNotes("Backfill guide-authoritative speaker notes from PPT_GUIDE.md after implementing the real slides.");
        }}

        async function main() {{
          const pptx = new pptxgen();
          pptx.layout = "LAYOUT_WIDE";
          pptx.author = "{author}";
          pptx.company = "{author}";
          pptx.subject = "{title}";
          pptx.title = "{title}";
          pptx.lang = "en-US";
          pptx.theme = {{
            headFontFace: "Aptos Display",
            bodyFontFace: "Aptos",
            lang: "en-US",
          }};
          buildCover(pptx);

          await pptx.writeFile({{ fileName: "deck.pptx" }});
        }}

        main().catch((error) => {{
          console.error(error);
          process.exitCode = 1;
        }});
        """
    )


def python_template(title: str, author: str) -> str:
    return textwrap.dedent(
        f"""\
        from __future__ import annotations

        from pathlib import Path

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        OUTPUT = Path(__file__).with_name("deck.pptx")
        PPT_GUIDE = Path(__file__).with_name("PPT_GUIDE.md")


        def placeholder_notes() -> list[str]:
            # Replace this with notes parsed from PPT_GUIDE.md.
            return ["Backfill guide-authoritative speaker notes from PPT_GUIDE.md after implementing the real slides."]


        def apply_speaker_notes(prs: Presentation, notes: list[str]) -> None:
            if len(prs.slides) != len(notes):
                raise ValueError(f"Slides ({{len(prs.slides)}}) and notes ({{len(notes)}}) count mismatch")
            for slide, note in zip(prs.slides, notes):
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.clear()
                notes_frame.text = note


        def build_cover(prs: Presentation) -> None:
            # Keep stable slide ids in function names, comments, and review notes.
            # Do not place ids like s01-cover on the visible slide unless explicitly requested.
            # If a slide needs an important formula, factor the OMML/equation helper out of the slide builder.
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(247, 245, 241)
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.2), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = "{title}"
            run = title_frame.paragraphs[0].runs[0]
            run.font.size = Pt(24)
            run.font.bold = True
            subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.4), Inches(0.8))
            subtitle_box.text_frame.text = "Replace this placeholder with audience-facing content from PPT_GUIDE.md."


        def build_presentation() -> Path:
            prs = Presentation()
            prs.slide_width = Inches(13.333333)
            prs.slide_height = Inches(7.5)
            prs.core_properties.author = "{author}"
            prs.core_properties.title = "{title}"
            prs.core_properties.subject = "{title}"
            build_cover(prs)
            prs.save(str(OUTPUT))
            # Save once before writing notes; python-pptx may not persist notes parts reliably on the first pass.
            notes_prs = Presentation(str(OUTPUT))
            apply_speaker_notes(notes_prs, placeholder_notes())
            notes_prs.save(str(OUTPUT))
            return OUTPUT


        if __name__ == "__main__":
            print(build_presentation())
        """
    )


def js_package_json() -> str:
    return textwrap.dedent(
        """\
        {
          "name": "deck-workspace",
          "private": true,
          "type": "commonjs",
          "scripts": {
            "build": "node generate_ppt.js"
          },
          "dependencies": {
            "pptxgenjs": "^4.0.1"
          }
        }
        """
    )


def python_requirements() -> str:
    return textwrap.dedent(
        """\
        python-pptx>=0.6.23
        PyMuPDF>=1.24.0
        Pillow>=10.0.0
        pdf2image>=1.16.0
        """
    )


def review_commands(generator_name: str, build_command: str, backend: str) -> str:
    sections = ["# Review Commands", ""]
    if backend == "python":
        sections.append(
            textwrap.dedent(
                """\
                ## Python Environment Setup

                ```bash
                python3 -m venv .venv
                source .venv/bin/activate
                pip install -r requirements.txt
                ```

                If `python-pptx` or related libraries are missing, try this local environment before falling back to JavaScript.
                """
            ).strip()
        )
        sections.append("")
    sections.append(
        textwrap.dedent(
            f"""\
            ## Build

            ```bash
            {build_command}
            ```

            ## Render To PDF

            ```bash
            soffice --headless --convert-to pdf --outdir rendered deck.pptx
            ```

            ## Render PDF To PNG

            ```bash
            pdftoppm -png rendered/deck.pdf rendered/slide
            ```

            ## Preferred One-Command Review Path

            ```bash
            python path/to/render_review.py deck.pptx --output-dir rendered
            ```

            ## Structural Validation

            ```bash
            python path/to/validate_deck.py deck.pptx --guide PPT_GUIDE.md --expect-notes --expect-slides <expected>
            ```

            Use this before handoff to catch wrong slide counts, missing notes parts, empty notes, or visible slide-id leakage.

            ## Workflow Reminder

            - Change structure, message, timing, or slide order in `PPT_GUIDE.md`.
            - Change layout, spacing, cropping, fonts, or visual fixes in `{generator_name}`.
            - Change formula choice, notes policy, or symbol explanations in `PPT_GUIDE.md`; change formula rendering or notes persistence bugs in `{generator_name}`.
            - When text overflows or wraps badly, use the ladder in `references/text-overflow.md` rather than shrinking fonts first.
            - Keep stable slide ids in the guide, code, and review notes rather than visible slide text unless explicitly requested.
            - If notes are part of delivery, verify the final `.pptx` still contains them and treat the post-notes file as the only handoff artifact.
            - If formulas are important, manually review the rendered formula pages after every meaningful edit.
            - Rebuild and rerender after every meaningful edit.
            """
        ).strip()
    )
    return "\n\n".join(sections) + "\n"


def review_template() -> str:
    return textwrap.dedent(
        """\
        # Review Notes

        ## Current Round

        - Scope: `TODO`
        - Reviewer: `TODO`
        - Render Used: `TODO`
        - Validator Run: `TODO` (e.g. `scripts/validate_deck.py deck.pptx --expect-notes`)

        ## Focus Checks

        - `Important formulas are visible when required, fit their containers, and explain key symbols`
        - `Speaker notes exist in the final deck and still match PPT_GUIDE.md when required`
        - `Summary/takeaway/closing-adjacent pages keep visible text in the audience's working language`
        - `High-risk text components (titles, chips, captions, number cards, summary cards) do not overflow or silently shrink`
        - `Highlight boxes and emphasis labels do not cover the content they annotate`
        - `Cropped figures and tables still match the guide-requested region; no stale assets reused`

        ## Open Issues

        - [ ] `slide-id` - `Describe the issue and route: guide/script/both`

        ## Closed Issues

        - None yet.

        ## Deferred

        - None.
        """
    )


def write_text(path: Path, content: str, force: bool) -> None:
    if path.exists() and not force:
        raise FileExistsError(f"{path} already exists; use --force to overwrite")
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def touch_keep(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.touch(exist_ok=True)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Initialize a guide-first deck workspace.")
    parser.add_argument("workspace", help="Path to the workspace directory to create.")
    parser.add_argument("--title", default="New Deck", help="Working title for the deck.")
    parser.add_argument("--author", default="Your Name", help="Presenter or owner name.")
    parser.add_argument("--audience", default="General audience", help="Intended audience.")
    parser.add_argument("--duration-minutes", type=int, default=20, help="Target duration in minutes.")
    parser.add_argument("--slides", type=int, default=10, help="Planned slide count.")
    parser.add_argument("--backend", choices=("js", "python"), default="js", help="Generator backend to scaffold.")
    parser.add_argument("--language", choices=("en", "zh"), default="en", help="Language for the guide template.")
    parser.add_argument("--force", action="store_true", help="Overwrite existing files in the workspace.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    workspace = Path(args.workspace).expanduser().resolve()
    workspace.mkdir(parents=True, exist_ok=True)

    generator_name = "generate_ppt.js" if args.backend == "js" else "generate_ppt.py"
    build_command = "node generate_ppt.js" if args.backend == "js" else "python generate_ppt.py"
    backend_name = "JavaScript / PptxGenJS" if args.backend == "js" else "Python / python-pptx"

    if workspace.parts[:2] == ("/", "tmp") or str(workspace).startswith("/tmp/"):
        print(
            "Warning: the workspace path looks transient. For maintainable deck work, place it inside the user's repo.",
            file=sys.stderr,
        )

    guide = english_guide(
        args.title,
        args.audience,
        args.duration_minutes,
        args.slides,
        args.author,
        generator_name,
        backend_name,
    )
    if args.language == "zh":
        guide = chinese_guide(
            args.title,
            args.audience,
            args.duration_minutes,
            args.slides,
            args.author,
            generator_name,
            backend_name,
        )

    write_text(workspace / "PPT_GUIDE.md", guide, args.force)
    if args.backend == "js":
        write_text(workspace / "generate_ppt.js", js_template(args.title, args.author), args.force)
        write_text(workspace / "package.json", js_package_json(), args.force)
    else:
        write_text(workspace / "generate_ppt.py", python_template(args.title, args.author), args.force)
        write_text(workspace / "requirements.txt", python_requirements(), args.force)
    write_text(workspace / "review" / "notes.md", review_template(), args.force)
    write_text(workspace / "review" / "commands.md", review_commands(generator_name, build_command, args.backend), args.force)
    touch_keep(workspace / "assets" / ".gitkeep")
    touch_keep(workspace / "rendered" / ".gitkeep")

    print(f"Initialized deck workspace at {workspace}")
    print("Created:")
    items = [workspace / "PPT_GUIDE.md"]
    if args.backend == "js":
        items.extend([workspace / "generate_ppt.js", workspace / "package.json"])
    else:
        items.extend([workspace / "generate_ppt.py", workspace / "requirements.txt"])
    items.extend(
        [
            workspace / "review" / "notes.md",
            workspace / "review" / "commands.md",
            workspace / "assets" / ".gitkeep",
            workspace / "rendered" / ".gitkeep",
        ]
    )
    for item in items:
        print(f"  - {item}")


if __name__ == "__main__":
    main()
