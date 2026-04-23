"""Self-demo deck for the deck-workflow skill.

Generates ``ClaudeCode_intro.pptx`` at the 16:9 wide layout using python-pptx.
Notes are written via the two-stage save pattern documented in
``deck-workflow/references/formulas-and-notes.md``.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt


HERE = Path(__file__).resolve().parent
OUTPUT = HERE / "ClaudeCode_intro.pptx"
PPT_GUIDE = HERE / "PPT_GUIDE.md"


NAVY = RGBColor(0x0F, 0x17, 0x2A)
NAVY_SOFT = RGBColor(0x1E, 0x29, 0x3B)
CORAL = RGBColor(0xFF, 0x6B, 0x35)
CORAL_SOFT = RGBColor(0xFF, 0xA4, 0x80)
CREAM = RGBColor(0xF5, 0xEF, 0xE0)
OFFWHITE = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x72, 0x80)
INK = RGBColor(0x1F, 0x23, 0x37)

FONT_SANS = "Ubuntu"
FONT_MONO = "DejaVu Sans Mono"

SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5


def inches_to_emu(value: float) -> int:
    return int(round(value * 914400))


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = inches_to_emu(SLIDE_W_IN)
    prs.slide_height = inches_to_emu(SLIDE_H_IN)


def fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape) -> None:
    shape.line.fill.background()


def set_background(slide, color: RGBColor) -> None:
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        inches_to_emu(SLIDE_W_IN),
        inches_to_emu(SLIDE_H_IN),
    )
    fill_solid(background, color)
    no_line(background)
    spTree = background._element.getparent()
    spTree.remove(background._element)
    spTree.insert(2, background._element)


def add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    color: RGBColor,
    rounded: bool = False,
) -> object:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        inches_to_emu(left),
        inches_to_emu(top),
        inches_to_emu(width),
        inches_to_emu(height),
    )
    if rounded:
        shape.adjustments[0] = 0.08
    fill_solid(shape, color)
    no_line(shape)
    return shape


def add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font: str = FONT_SANS,
    size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    letter_spacing: float | None = None,
) -> object:
    box = slide.shapes.add_textbox(
        inches_to_emu(left),
        inches_to_emu(top),
        inches_to_emu(width),
        inches_to_emu(height),
    )
    frame = box.text_frame
    frame.margin_left = Emu(0)
    frame.margin_right = Emu(0)
    frame.margin_top = Emu(0)
    frame.margin_bottom = Emu(0)
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing is not None:
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))
    return box


def add_multiline(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[tuple[str, dict]],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: float = 1.15,
) -> object:
    box = slide.shapes.add_textbox(
        inches_to_emu(left),
        inches_to_emu(top),
        inches_to_emu(width),
        inches_to_emu(height),
    )
    frame = box.text_frame
    frame.margin_left = Emu(0)
    frame.margin_right = Emu(0)
    frame.margin_top = Emu(0)
    frame.margin_bottom = Emu(0)
    frame.word_wrap = True
    for index, (text, opts) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = opts.get("align", align)
        paragraph.line_spacing = line_spacing
        if opts.get("space_before") is not None:
            paragraph.space_before = Pt(opts["space_before"])
        run = paragraph.add_run()
        run.text = text
        run.font.name = opts.get("font", FONT_SANS)
        run.font.size = Pt(opts.get("size", 14))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", INK)
    return box


def add_circle(slide, cx: float, cy: float, diameter: float, color: RGBColor) -> object:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        inches_to_emu(cx - diameter / 2),
        inches_to_emu(cy - diameter / 2),
        inches_to_emu(diameter),
        inches_to_emu(diameter),
    )
    fill_solid(shape, color)
    no_line(shape)
    return shape


def add_dot_grid(slide, left: float, top: float, cols: int, rows: int, gap: float, diameter: float, color: RGBColor) -> None:
    for r in range(rows):
        for c in range(cols):
            cx = left + c * gap
            cy = top + r * gap
            add_circle(slide, cx, cy, diameter, color)


def kicker(slide, text: str, left: float, top: float, color: RGBColor = CORAL, width: float = 11.5) -> None:
    add_text(slide, left, top, width, 0.4, text, size=11, bold=True, color=color, letter_spacing=2)


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)

    add_rect(slide, 0.6, 0.55, 0.45, 0.12, CORAL)
    kicker(slide, "AN INTRO FOR HUMANS WHO NEVER WROTE CODE", 1.15, 0.48, color=CORAL, width=10.5)

    add_multiline(
        slide,
        0.6, 1.9, 11.8, 3.2,
        [
            ("Meet Claude Code.", {"size": 72, "bold": True, "color": WHITE, "font": FONT_SANS}),
            ("Your patient AI coworker", {"size": 34, "color": CREAM, "font": FONT_SANS, "space_before": 18}),
            ("who never sighs.", {"size": 34, "color": CORAL_SOFT, "bold": True, "font": FONT_SANS}),
        ],
        line_spacing=1.05,
    )

    add_dot_grid(slide, 11.5, 0.55, cols=6, rows=4, gap=0.22, diameter=0.08, color=CORAL_SOFT)

    add_text(
        slide, 0.6, 6.75, 12.1, 0.35,
        "6 slides.  60 seconds.  No prior experience required.",
        size=13, color=CREAM,
    )
    add_text(
        slide, 0.6, 7.05, 6.0, 0.3,
        "github.com/HansBug/deck-workflow-skill",
        size=11, color=CORAL_SOFT, font=FONT_MONO,
    )


def slide_what_is_it(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)

    kicker(slide, "WHAT EVEN IS IT", 0.75, 0.6, color=CORAL)

    add_multiline(
        slide,
        0.75, 1.0, 11.8, 1.8,
        [
            ("It is not a chatbot that writes code.", {"size": 30, "bold": True, "color": NAVY}),
            ("It is a coworker that does the work.", {"size": 30, "bold": True, "color": CORAL}),
        ],
        line_spacing=1.15,
    )

    cards = [
        ("A senior engineer", "Who has infinite patience and zero ego."),
        ("Living in your terminal", "You type, it types back, it runs the thing."),
        ("Explains like you want", 'Ask for "like I am five" and it means it.'),
    ]

    card_w = 3.85
    card_h = 3.0
    gap = 0.25
    total = 3 * card_w + 2 * gap
    start_left = (SLIDE_W_IN - total) / 2
    card_top = 3.55

    for index, (title, body) in enumerate(cards):
        left = start_left + index * (card_w + gap)
        add_rect(slide, left, card_top, card_w, card_h, WHITE)
        add_rect(slide, left, card_top, card_w, 0.08, CORAL)
        add_text(
            slide, left + 0.3, card_top + 0.35, card_w - 0.6, 0.3,
            f"0{index + 1}",
            size=14, bold=True, color=CORAL, letter_spacing=2,
        )
        add_text(
            slide, left + 0.3, card_top + 0.8, card_w - 0.6, 0.8,
            title,
            size=22, bold=True, color=NAVY,
        )
        add_text(
            slide, left + 0.3, card_top + 1.55, card_w - 0.6, 1.1,
            body,
            size=14, color=INK,
        )

    add_text(
        slide, 0.75, 6.95, 12.0, 0.3,
        "s02  ·  Pick whichever metaphor calms your nerves.",
        size=11, italic=True, color=MUTED,
    )


def slide_how_it_feels(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)

    kicker(slide, "HOW IT FEELS", 0.75, 0.6, color=CORAL)

    add_multiline(
        slide,
        0.75, 1.0, 11.8, 1.4,
        [
            ("You just... talk to it.", {"size": 32, "bold": True, "color": NAVY}),
            ("A real conversation. Not a feature list.", {"size": 16, "color": MUTED, "space_before": 4}),
        ],
        line_spacing=1.15,
    )

    term_left = 1.25
    term_top = 2.75
    term_w = 10.85
    term_h = 3.9
    add_rect(slide, term_left, term_top, term_w, term_h, NAVY, rounded=True)

    dot_y = term_top + 0.35
    for i, dot_color in enumerate([CORAL, CORAL_SOFT, CREAM]):
        add_circle(slide, term_left + 0.45 + i * 0.3, dot_y, 0.18, dot_color)

    add_text(
        slide, term_left + 1.6, term_top + 0.23, 6.0, 0.28,
        "claude — interactive shell",
        size=11, color=CREAM, font=FONT_MONO,
    )

    dialog = [
        ("user",  "hey can you build me a website that tells bad dad jokes"),
        ("claude", "Sure. Setting up a single page with a button that pulls jokes from an API."),
        ("claude", "Creating index.html, style.css, and script.js now."),
        ("user",  "looks great. can you make the background neon pink"),
        ("claude", "Done. It is extremely pink now. Regrettable but functional."),
    ]

    lines: list[tuple[str, dict]] = []
    for speaker, text in dialog:
        if speaker == "user":
            lines.append(("$ " + text, {"font": FONT_MONO, "size": 14, "color": CORAL_SOFT, "bold": True}))
        else:
            lines.append(("  " + text, {"font": FONT_MONO, "size": 14, "color": CREAM}))

    add_multiline(
        slide,
        term_left + 0.5, term_top + 0.85, term_w - 1.0, term_h - 1.1,
        lines,
        line_spacing=1.5,
    )

    add_text(
        slide, 0.75, 6.9, 12.0, 0.35,
        "That is it. That is the interaction.",
        size=14, italic=True, color=NAVY, align=PP_ALIGN.CENTER,
    )


def slide_why_beginners_love_it(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)

    kicker(slide, "WHY IT STICKS", 0.75, 0.6, color=CORAL)

    add_text(
        slide, 0.75, 1.0, 11.8, 0.7,
        "Why beginners actually stick with it.",
        size=30, bold=True, color=NAVY,
    )
    add_text(
        slide, 0.75, 1.7, 11.8, 0.4,
        "Three things that make the learning curve flatter than you expect.",
        size=15, color=MUTED,
    )

    cards = [
        ("Speaks human", "No jargon unless you ask for it. No condescension either."),
        ("Shows its work", "Watch it think, plan, and edit files in real time."),
        ("Fixes its own bugs", "If it breaks the build, it usually fixes it before you notice."),
    ]

    card_w = 3.85
    card_h = 3.85
    gap = 0.25
    total = 3 * card_w + 2 * gap
    start_left = (SLIDE_W_IN - total) / 2
    card_top = 2.8

    for index, (title, body) in enumerate(cards):
        left = start_left + index * (card_w + gap)
        add_rect(slide, left, card_top, card_w, card_h, WHITE)
        add_text(
            slide, left + 0.35, card_top + 0.3, card_w - 0.7, 1.4,
            f"0{index + 1}",
            size=64, bold=True, color=CORAL,
        )
        add_text(
            slide, left + 0.35, card_top + 1.85, card_w - 0.7, 0.55,
            title,
            size=22, bold=True, color=NAVY,
        )
        add_text(
            slide, left + 0.35, card_top + 2.55, card_w - 0.7, 1.1,
            body,
            size=13, color=INK,
        )

    add_text(
        slide, 0.75, 6.95, 12.0, 0.3,
        "s04  ·  No one learns when they feel stupid.",
        size=11, italic=True, color=MUTED,
    )


def slide_what_you_can_build(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)

    kicker(slide, "WHAT YOU CAN BUILD", 0.75, 0.6, color=CORAL)

    add_text(
        slide, 0.75, 1.0, 11.8, 0.7,
        "Things you can ship on day one.",
        size=30, bold=True, color=NAVY,
    )
    add_text(
        slide, 0.75, 1.7, 11.8, 0.4,
        "Zero experience, one laptop, a cup of coffee.",
        size=15, color=MUTED,
    )

    tiles = [
        ("A personal website", "~ 5 minutes"),
        ("A script that renames 3000 files", "~ 3 minutes"),
        ("A translator for your grandma's letters", "~ 10 minutes"),
        ("A bot that books your climbing gym at 7 AM", "~ 20 minutes, lightly unethical"),
    ]

    tile_w = 5.6
    tile_h = 2.05
    gap_x = 0.3
    gap_y = 0.3
    total_w = 2 * tile_w + gap_x
    start_left = (SLIDE_W_IN - total_w) / 2
    start_top = 2.65

    for index, (title, meta) in enumerate(tiles):
        row = index // 2
        col = index % 2
        left = start_left + col * (tile_w + gap_x)
        top = start_top + row * (tile_h + gap_y)
        add_rect(slide, left, top, tile_w, tile_h, WHITE)
        add_rect(slide, left, top, 0.08, tile_h, CORAL)
        add_text(
            slide, left + 0.35, top + 0.35, tile_w - 0.7, 1.1,
            title,
            size=22, bold=True, color=NAVY,
        )
        add_text(
            slide, left + 0.35, top + 1.45, tile_w - 0.7, 0.45,
            meta,
            size=13, italic=True, color=CORAL,
            font=FONT_MONO,
        )

    add_text(
        slide, 0.75, 7.0, 12.0, 0.3,
        "s05  ·  The gym bot slide always gets the laugh.",
        size=11, italic=True, color=MUTED,
    )


def slide_get_started(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)

    add_rect(slide, 0.6, 0.55, 0.45, 0.12, CORAL)
    kicker(slide, "GET STARTED", 1.15, 0.48, color=CORAL)

    add_multiline(
        slide,
        0.6, 1.1, 11.8, 1.6,
        [
            ("Go build something.", {"size": 58, "bold": True, "color": WHITE}),
            ("Two commands. That is the entire install.", {"size": 20, "color": CREAM, "space_before": 6}),
        ],
        line_spacing=1.1,
    )

    cards = [
        ("STEP 1 — INSTALL ONCE", "npm install -g @anthropic-ai/claude-code", 7.8),
        ("STEP 2 — RUN", "claude", 3.5),
    ]

    card_h = 2.2
    gap = 0.4
    total_w = cards[0][2] + cards[1][2] + gap
    start_left = (SLIDE_W_IN - total_w) / 2
    card_top = 3.6

    left_cursor = start_left
    for title, code, card_w in cards:
        add_rect(slide, left_cursor, card_top, card_w, card_h, NAVY_SOFT, rounded=True)
        add_rect(slide, left_cursor + 0.35, card_top + 0.4, 0.3, 0.08, CORAL)
        add_text(
            slide, left_cursor + 0.75, card_top + 0.32, card_w - 1.1, 0.3,
            title, size=11, bold=True, color=CORAL_SOFT, letter_spacing=2,
        )
        add_text(
            slide, left_cursor + 0.35, card_top + 1.0, card_w - 0.7, 0.9,
            code, size=18, bold=True, color=CREAM, font=FONT_MONO,
        )
        left_cursor += card_w + gap

    add_text(
        slide, 0.6, 6.6, 12.1, 0.35,
        "claude.com/claude-code",
        size=14, color=CORAL_SOFT, font=FONT_MONO, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, 0.6, 6.95, 12.1, 0.35,
        "Now close this deck. Go build.",
        size=16, italic=True, color=CREAM, align=PP_ALIGN.CENTER,
    )


SLIDE_BUILDERS: list[Callable[[Presentation], None]] = [
    slide_cover,
    slide_what_is_it,
    slide_how_it_feels,
    slide_why_beginners_love_it,
    slide_what_you_can_build,
    slide_get_started,
]


GUIDE_SLIDE_HEADING = re.compile(r"^###\s+(s\d{2,3}-[^\s]+)", re.MULTILINE)
READ_ALOUD_FIELD = "Speaker Notes:"


@dataclass
class GuideBlock:
    slide_id: str
    block: str


def parse_guide_blocks(guide_path: Path) -> list[GuideBlock]:
    text = guide_path.read_text(encoding="utf-8")
    headings = list(GUIDE_SLIDE_HEADING.finditer(text))
    blocks: list[GuideBlock] = []
    for index, match in enumerate(headings):
        start = match.end()
        end = headings[index + 1].start() if index + 1 < len(headings) else len(text)
        blocks.append(GuideBlock(slide_id=match.group(1), block=text[start:end]))
    return blocks


def extract_read_aloud(block: str) -> str:
    lines = block.splitlines()
    buffer: list[str] = []
    collecting = False
    for line in lines:
        stripped = line.strip()
        if not collecting:
            if stripped.startswith("- ") and READ_ALOUD_FIELD in stripped:
                collecting = True
            continue
        if not stripped:
            if buffer:
                break
            continue
        if stripped.startswith("- ") and stripped.endswith(":") and buffer:
            break
        if stripped.startswith("- "):
            buffer.append(stripped[2:].strip(" `"))
        else:
            buffer.append(stripped.strip(" `"))
    return " ".join(entry for entry in buffer if entry)


def extract_notes_from_guide(guide_path: Path) -> list[str]:
    blocks = parse_guide_blocks(guide_path)
    notes: list[str] = []
    for block in blocks:
        line = extract_read_aloud(block.block)
        if not line:
            raise RuntimeError(f"Guide block {block.slide_id} has no speaker notes")
        notes.append(line)
    return notes


def apply_speaker_notes(prs: Presentation, notes: list[str]) -> None:
    if len(prs.slides) != len(notes):
        raise ValueError(
            f"Slide count {len(prs.slides)} does not match notes count {len(notes)}"
        )
    for slide, note in zip(prs.slides, notes):
        frame = slide.notes_slide.notes_text_frame
        frame.clear()
        frame.text = note


def build_presentation() -> Path:
    notes = extract_notes_from_guide(PPT_GUIDE)
    if len(notes) != len(SLIDE_BUILDERS):
        raise RuntimeError(
            f"Guide has {len(notes)} slide blocks but generator has {len(SLIDE_BUILDERS)} builders"
        )

    prs = Presentation()
    set_slide_size(prs)
    for builder in SLIDE_BUILDERS:
        builder(prs)
    prs.save(str(OUTPUT))

    notes_prs = Presentation(str(OUTPUT))
    apply_speaker_notes(notes_prs, notes)
    notes_prs.save(str(OUTPUT))

    return OUTPUT


def main() -> None:
    path = build_presentation()
    print(f"Deck written to: {path}")
    print(f"Slides: {len(SLIDE_BUILDERS)}")


if __name__ == "__main__":
    main()
